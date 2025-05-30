import pytest
from pathlib import Path
import tempfile
from PyPDF2 import PdfWriter
from pptx import Presentation
from pptx.util import Inches
from src.presentation_processor import PresentationProcessor
from src.content_analyzer import ContentAnalyzer, SlideContent

@pytest.fixture
def processor():
    return PresentationProcessor()

@pytest.fixture
def temp_dir():
    with tempfile.TemporaryDirectory() as tmpdirname:
        yield Path(tmpdirname)

@pytest.fixture
def sample_pdf(temp_dir):
    pdf_path = temp_dir / "test.pdf"
    
    # Create a simple PDF with two pages
    writer = PdfWriter()
    
    # First page
    page1 = writer.add_blank_page(width=612, height=792)
    page1.merge_page(PdfWriter().add_blank_page(width=612, height=792))
    page1.merge_page(PdfWriter().add_blank_page(width=612, height=792))
    
    # Second page
    page2 = writer.add_blank_page(width=612, height=792)
    page2.merge_page(PdfWriter().add_blank_page(width=612, height=792))
    
    with open(pdf_path, 'wb') as f:
        writer.write(f)
    
    return pdf_path

@pytest.fixture
def sample_pptx(temp_dir):
    pptx_path = temp_dir / "test.pptx"
    
    # Create a simple PPTX with two slides
    prs = Presentation()
    
    # First slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    title1.text = "First Slide"
    
    # Add some text
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1)
    textbox = slide1.shapes.add_textbox(left, top, width, height)
    textbox.text = "This is the first slide content"
    
    # Second slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[0])
    title2 = slide2.shapes.title
    title2.text = "Second Slide"
    
    # Add some text
    textbox2 = slide2.shapes.add_textbox(left, top, width, height)
    textbox2.text = "This is the second slide content"
    
    prs.save(pptx_path)
    return pptx_path

@pytest.fixture
def content_analyzer():
    return ContentAnalyzer()

@pytest.fixture
def sample_slides():
    return [
        SlideContent(
            slide_number=1,
            text="Title Slide\n• First point\n• Second point",
            images=[],
            metadata={"type": "title"}
        ),
        SlideContent(
            slide_number=2,
            text="Content Slide\n• Main point 1\n• Main point 2\n• Main point 3",
            images=[],
            metadata={"type": "content"}
        ),
        SlideContent(
            slide_number=3,
            text="Image Slide with some text",
            images=["image1.jpg"],
            metadata={"type": "image"}
        )
    ]
