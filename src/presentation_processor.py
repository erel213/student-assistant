from pathlib import Path
from typing import Dict, List, Optional, Union
import PyPDF2
from pptx import Presentation
from dataclasses import dataclass
import logging

@dataclass
class SlideContent:
    """Represents the content extracted from a single slide."""
    slide_number: int
    text: str
    images: List[Dict[str, str]]
    metadata: Dict[str, str]

class PresentationProcessor:
    """Processes presentation files (PDF and PPTX) to extract content and structure."""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    def process_file(self, file_path: Union[str, Path]) -> List[SlideContent]:
        """
        Process a presentation file and extract its content.
        
        Args:
            file_path: Path to the presentation file (PDF or PPTX)
            
        Returns:
            List of SlideContent objects containing extracted information
            
        Raises:
            ValueError: If file format is not supported
            FileNotFoundError: If file does not exist
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
            
        if file_path.suffix.lower() == '.pdf':
            return self._process_pdf(file_path)
        elif file_path.suffix.lower() == '.pptx':
            return self._process_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_path.suffix}")
    
    def _process_pdf(self, file_path: Path) -> List[SlideContent]:
        """Process a PDF presentation file."""
        slides = []
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    # Create slide content
                    slide = SlideContent(
                        slide_number=page_num + 1,
                        text=text,
                        images=[],  # PDF image extraction would require additional processing
                        metadata={
                            'page_count': len(pdf_reader.pages),
                            'file_type': 'pdf'
                        }
                    )
                    slides.append(slide)
                    
        except Exception as e:
            self.logger.error(f"Error processing PDF file: {e}")
            raise
            
        return slides
    
    def _process_pptx(self, file_path: Path) -> List[SlideContent]:
        """Process a PPTX presentation file."""
        slides = []
        
        try:
            presentation = Presentation(file_path)
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                # Extract text from all shapes
                text_content = []
                images = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
                    
                    # Handle images
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        image_info = {
                            'type': 'image',
                            'position': {
                                'left': shape.left,
                                'top': shape.top,
                                'width': shape.width,
                                'height': shape.height
                            }
                        }
                        images.append(image_info)
                
                # Create slide content
                slide_content = SlideContent(
                    slide_number=slide_num,
                    text='\n'.join(text_content),
                    images=images,
                    metadata={
                        'slide_count': len(presentation.slides),
                        'file_type': 'pptx'
                    }
                )
                slides.append(slide_content)
                
        except Exception as e:
            self.logger.error(f"Error processing PPTX file: {e}")
            raise
            
        return slides
